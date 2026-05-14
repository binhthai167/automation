from sqlalchemy import create_engine, text
import pandas as pd
from datetime import datetime
from openpyxl import load_workbook
import calendar
from copy import copy
import xlwings as xw
import os
from PIL import ImageGrab
from pptx import Presentation
from pptx.util import Inches

engine = create_engine(
    "mysql+pymysql://admin:Nidec%40123456@192.168.150.210:3306/edms"
)

query = """
SELECT 
    d.name AS department,

    -- tổng form từ trước tới nay
    COUNT(DISTINCT fd.id) AS total_forms_all_time,

    -- tổng form tháng này
    COUNT(DISTINCT CASE 
        WHEN fd.date_create >= DATE_FORMAT(CURDATE(),'%%Y-%%m-01')
        AND fd.date_create < DATE_FORMAT(CURDATE() + INTERVAL 1 MONTH,'%%Y-%%m-01')
        THEN fd.id 
    END) AS total_forms_this_month

FROM edms.wf_form_data fd
JOIN edms.wf_form f ON fd.form_id = f.id
JOIN edms.form_of_folder ff ON f.id = ff.form_id
JOIN edms.folder_of_dept fod ON ff.folder_id = fod.id
JOIN edms.org_department d ON fod.dept_id = d.id

WHERE f.parent_id IS NULL
AND d.name <> 'Draft'
AND d.name <> 'PRESS'
AND d.name <> 'LEAD_CREW'

GROUP BY d.name;
"""

df = pd.read_sql(query, engine)
print(df)

file_path = "report/List report EDMS 28022026.xlsx"
app = xw.App(visible=False)
wb = app.books.open(file_path)
ws = wb.sheets[0]

SUM_COL = 6
HEADER_ROW =3
DATA_START_ROW = 4
now = datetime.now()
header_cell = ws.range((HEADER_ROW, SUM_COL-1)).value
is_current_month = False
if header_cell:
    try:
        header_date = datetime.strptime(str(header_cell), "%m/%d/%Y")
        if header_date.month == now.month and header_date.year == now.year:
            is_current_month = True
    except:
        pass
if not is_current_month:
    # xóa tháng cũ
    ws.range("C:C").delete()

    # thêm cột tháng mới
    ws.api.Columns(SUM_COL-1).Insert()
    last_day = calendar.monthrange(now.year, now.month)[1]

    header = f"{now.month}/{now.day}/{now.year}"
    ws.range((HEADER_ROW, SUM_COL-1)).value = header
    print("thêm cột tháng mới")
else:
    print("Tháng hiện tại đã tồn tại, chỉ cập nhật dữ liệu")

# -------------------------
# điền dữ liệu tháng mới
# -------------------------
last_row = 8

dept_list = ws.range((DATA_START_ROW, 2), (last_row, 2)).value

# -------------------------
# dữ liệu từ DB
# -------------------------

dept_month = dict(zip(df["department"], df["total_forms_this_month"]))
dept_total = dict(zip(df["department"], df["total_forms_all_time"]))

month_values = []
sum_values = []
for dept in dept_list:

    if dept == "Total":
        month_values.append(None)
        sum_values.append(None)
        continue

    month_values.append(dept_month.get(dept, 0))
    sum_values.append(dept_total.get(dept, 0))

total_month = sum(v for v in month_values if isinstance(v, (int, float)))
total_sum = sum(v for v in sum_values if isinstance(v, (int, float)))

for i, dept in enumerate(dept_list):

    if dept == "Total":
        month_values[i] = total_month
        sum_values[i] = total_sum
        break

# -------------------------
# ghi dữ liệu vào Excel (batch)
# -------------------------

ws.range((DATA_START_ROW, SUM_COL-1)).options(transpose=True).value = month_values
ws.range((DATA_START_ROW, SUM_COL)).options(transpose=True).value = sum_values

chart = ws.api.ChartObjects(1).Chart

chart.SetSourceData(
    ws.range((3,2), (8,5)).api
)

chart.PlotBy = 1


rng = ws.range("H3:U23")

rng.api.CopyPicture(Appearance=1, Format=2)

img = ImageGrab.grabclipboard()

img_path = "reportEDMS.png"
img.save(img_path)

wb.save()
wb.close()
app.quit()

ppt_path = "report/3月10日一般経費報告 (IT) .pptx"

prs = Presentation(ppt_path)

slide = prs.slides[0]   # slide 1
target_shape =  None
for shape in slide.shapes:
    if shape.name == "Chart 1":
        target_shape = shape
        break
if target_shape is None:
    raise Exception("Không tìm thấy chart trong slide")
left = target_shape.left
top = target_shape.top
width = target_shape.width
height = target_shape.height

# xóa ảnh cũ
sp = target_shape._element
sp.getparent().remove(sp)

# chèn ảnh mới
pic = slide.shapes.add_picture(
    img_path,
    left,
    top,
    width=width,
    height=height
)
pic.name = "Chart 1"
# prs.save(f"report/{now.month}月{now.day}日一般経費報告 (IT).pptx")
prs.save(ppt_path)
print("Report updated successfully")


