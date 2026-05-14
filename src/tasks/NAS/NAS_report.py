import pandas as pd
from openpyxl import load_workbook
import subprocess
import xlwings as xw
from PIL import ImageGrab
import time
from pptx import Presentation
from datetime import datetime

# subprocess.run(["scan_nas.bat"])

encodings = ["utf-8-sig", "utf-8", "cp1258", "cp1252", "mbcs", "cp850"]

for enc in encodings:
    try:
        df = pd.read_csv("report.csv", encoding=enc)
        print("OK:", enc)
        print(df.head())
        break
    except Exception as e:
        print(enc, "fail")

# # bỏ khoảng trắng ở header
df.columns = df.columns.str.strip()

# # chuyển cột Size thành số
# df["Size"] = pd.to_numeric(df["Size"], errors="coerce")

# đổi sang GB
df["Size(GB)"] = (df["SizeBytes"] / (1024**3)).round(0)

# # # nếu muốn bỏ cột Size cũ
df = df.drop(columns=["SizeBytes"])

total_size = df["Size(GB)"].sum()
df["Percentage"] = (df["Size(GB)"] / total_size * 100).round(0)
df.insert(0, "No", range(1, len(df) + 1))
file = "report.xlsx"

book = load_workbook(file)
sheet = book["Sheet1"]

start_row = 2
for i, row in df.iterrows():
    sheet.cell(row=start_row+i, column=1).value = row["No"]
    sheet.cell(row=start_row+i, column=2).value = row["Folder"]
    sheet.cell(row=start_row+i, column=3).value = row["FileCount"]
    sheet.cell(row=start_row+i, column=4).value = row["Size(GB)"]
    sheet.cell(row=start_row+i, column=5).value = row["Percentage"]
    cell = sheet.cell(row=start_row+i, column=5)
    cell.value = row["Percentage"] / 100
    cell.number_format = '0%'
book.save(file)

wb = xw.Book(file)
ws = wb.sheets["Sheet1"]

def capture(range_str, filename):

    rng = ws.range(range_str)

    rng.api.CopyPicture(Appearance=1, Format=2)

    time.sleep(1)

    img = ImageGrab.grabclipboard()
    img.save(filename)
# ảnh 1
capture("A1:E17", "snapshot/table.png")

# ảnh 2
capture("G3:Q19", "snapshot/chart.png")
wb.close()

ppt_path = r"D:\Automation\src\tasks\EDMS\report\3月10日一般経費報告 (IT) .pptx"
img1 = "snapshot/table.png"
img2 = "snapshot/chart.png"

prs = Presentation(ppt_path)
# slide 2 (index bắt đầu từ 0)
slide = prs.slides[1]


def replace_image(slide, shape_name, image_path):

    target_shape = None

    for shape in slide.shapes:
        if shape.name == shape_name:
            target_shape = shape
            break

    if target_shape is None:
        raise Exception(f"Không tìm thấy {shape_name}")

    # lưu vị trí
    left = target_shape.left
    top = target_shape.top
    width = target_shape.width
    height = target_shape.height

    # xóa shape cũ
    sp = target_shape._element
    sp.getparent().remove(sp)

    # chèn ảnh mới
    pic = slide.shapes.add_picture(
        image_path,
        left,
        top,
        width=width,
        height=height
    )

    pic.name = shape_name


# thay 2 ảnh
replace_image(slide, "Picture 1", img1)
replace_image(slide, "Picture 3", img2)
now = datetime.now()
prs.save(f"report/{now.month}月{now.day}日一般経費報告 (IT).pptx")
print("done")